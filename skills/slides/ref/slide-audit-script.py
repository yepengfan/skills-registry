#!/usr/bin/env python3
"""Comprehensive slide audit script — checks all CI&T design rules.

Usage: python3 slide-audit-script.py Harness_Engineering_CIT.pptx
"""

import sys
from pptx import Presentation

# Brand colors
NAVY = "000050"
CORAL = "FA5A50"
MAROON = "690037"
LTBLUE = "B4DCFA"
LTPURPLE = "FAB9FF"
WHITE = "FFFFFF"

LIGHT_FILLS = {WHITE, LTBLUE, LTPURPLE, "F5F5F5", "F0F0F2"}
DARK_FILLS = {NAVY, MAROON}


def audit(pptx_path):
    prs = Presentation(pptx_path)
    violations = []

    for si, slide in enumerate(prs.slides):
        slide_bg = "?"
        try:
            slide_bg = str(slide.background.fill.fore_color.rgb)
        except:
            pass

        is_light_bg = slide_bg in LIGHT_FILLS
        is_coral_bg = slide_bg == CORAL

        # Collect all filled shapes
        filled = []
        for shape in slide.shapes:
            try:
                fc = str(shape.fill.fore_color.rgb)
                filled.append({
                    'c': fc, 'l': shape.left, 't': shape.top,
                    'r': shape.left + shape.width, 'b': shape.top + shape.height,
                    'a': shape.width * shape.height,
                    'w_pt': shape.width / 12700, 'h_pt': shape.height / 12700
                })
            except:
                pass

        filled_asc = sorted(filled, key=lambda f: f['a'])

        # ── Check 1: Same-color nesting ──────────────────────────────────
        for child in filled:
            for parent in filled:
                if child is parent or child['a'] >= parent['a']:
                    continue
                if (parent['l'] <= child['l'] and child['r'] <= parent['r'] and
                    parent['t'] <= child['t'] and child['b'] <= parent['b']):
                    if child['c'] == parent['c'] and child['c'] not in (WHITE, NAVY):
                        violations.append(
                            f"S{si+1}: NEST #{child['c']} inside #{parent['c']}")
                        break

        # ── Check 2: LTBLUE components on LTBLUE bg ─────────────────────
        if slide_bg == LTBLUE:
            coral_rects = [(f['l'], f['t'], f['r'], f['b'])
                          for f in filled if f['c'] == CORAL]
            for f in filled:
                if f['c'] == LTBLUE and f['a'] > 914400 * 914400 * 0.02:
                    in_coral = any(
                        cl <= f['l'] and f['r'] <= cr and ct <= f['t'] and f['b'] <= cb
                        for cl, ct, cr, cb in coral_rects)
                    if not in_coral:
                        violations.append(f"S{si+1}: LTBLUE component on LTBLUE bg")

        # ── Check 3: CORAL components on CORAL bg ────────────────────────
        if slide_bg == CORAL:
            white_rects = [(f['l'], f['t'], f['r'], f['b'])
                          for f in filled if f['c'] == WHITE]
            for f in filled:
                if f['c'] == CORAL and f['a'] > 914400 * 914400 * 0.02:
                    in_white = any(
                        wl <= f['l'] and f['r'] <= wr and wt <= f['t'] and f['b'] <= wb
                        for wl, wt, wr, wb in white_rects)
                    if not in_white:
                        violations.append(f"S{si+1}: CORAL component on CORAL bg")

        # ── Check 4: Text color violations ───────────────────────────────
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if not para.text.strip():
                    continue

                # Get color from paragraph or run level
                tc = None
                try:
                    if para.font.color and para.font.color.rgb:
                        tc = str(para.font.color.rgb)
                except:
                    pass
                if tc is None:
                    for run in para.runs:
                        try:
                            if run.font.color and run.font.color.rgb:
                                tc = str(run.font.color.rgb)
                                break
                        except:
                            pass
                if tc is None:
                    continue

                # Find innermost background
                sx, sy = shape.left, shape.top
                ebg = slide_bg
                for f in filled_asc:
                    if f['l'] <= sx <= f['r'] and f['t'] <= sy <= f['b']:
                        ebg = f['c']

                txt = para.text[:25]

                # Same color as background
                if tc == ebg:
                    violations.append(
                        f"S{si+1}: #{tc} text on #{ebg} (SAME): '{txt}'")

                # WHITE text on light background
                if tc == WHITE and ebg in LIGHT_FILLS:
                    violations.append(
                        f"S{si+1}: WHITE text on light #{ebg}: '{txt}'")

                # Light text on light background
                if tc in (LTBLUE, LTPURPLE) and ebg in LIGHT_FILLS:
                    violations.append(
                        f"S{si+1}: light #{tc} on light #{ebg}: '{txt}'")

        # ── Check 5: Accent bar violations ───────────────────────────────
        for f in filled:
            is_thin = (f['w_pt'] <= 5 or f['h_pt'] <= 5) and f['w_pt'] > 0.5 and f['h_pt'] > 0.5
            if not is_thin:
                continue

            # Light bars on light bg
            if is_light_bg and f['c'] in LIGHT_FILLS:
                violations.append(
                    f"S{si+1}: LIGHT bar #{f['c']} on LIGHT bg #{slide_bg}")

            # Dark bars on dark bg (CORAL)
            if is_coral_bg and f['c'] in DARK_FILLS:
                violations.append(
                    f"S{si+1}: DARK bar #{f['c']} on CORAL bg")

            # Bar same as slide bg
            if f['c'] == slide_bg:
                violations.append(
                    f"S{si+1}: bar #{f['c']} = slide bg (invisible)")

    return violations


def main():
    path = sys.argv[1] if len(sys.argv) > 1 else "Harness_Engineering_CIT.pptx"
    violations = audit(path)

    # Known false positives: S11 WHITE text inside CORAL fnode on WHITE swim lane
    violations = [v for v in violations if not (
        "S11:" in v and "WHITE text on light #FFFFFF" in v)]

    print(f"\n{'='*60}")
    print(f"AUDIT RESULTS: {path}")
    print(f"{'='*60}")
    print(f"Violations: {len(violations)}")

    if violations:
        for v in violations:
            print(f"  ❌ {v}")
    else:
        print("  ✅ ALL CLEAR — no violations found!")

    print(f"{'='*60}\n")
    return len(violations)


if __name__ == '__main__':
    sys.exit(main())
